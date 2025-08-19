from .storage import download_bytes
import mimetypes

try:
    import docx
except Exception:
    docx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import pptx
except Exception:
    pptx = None

def _extract_docx_to_html(bytes_data: bytes) -> bytes:
    if not docx:
        raise RuntimeError('python-docx not installed')
    from io import BytesIO
    doc = docx.Document(BytesIO(bytes_data))
    parts = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(f"<p>{text}</p>")
    return ('\n'.join(parts)).encode('utf-8')

def _extract_xlsx_to_html(bytes_data: bytes) -> bytes:
    if not openpyxl:
        raise RuntimeError('openpyxl not installed')
    from io import BytesIO
    wb = openpyxl.load_workbook(filename=BytesIO(bytes_data), read_only=True, data_only=True)
    sheet = wb[wb.sheetnames[0]]
    rows = []
    for r in sheet.iter_rows(values_only=True):
        cells = ''.join(f'<td>{(c if c is not None else "")}</td>' for c in r)
        rows.append(f'<tr>{cells}</tr>')
    html = '<table border="1">' + '\n'.join(rows) + '</table>'
    return html.encode('utf-8')

def _extract_pptx_to_html(bytes_data: bytes) -> bytes:
    if not pptx:
        raise RuntimeError('python-pptx not installed')
    from io import BytesIO
    prs = pptx.Presentation(BytesIO(bytes_data))
    slides_html = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides_html.append(f"<h3>Slide {i}</h3><p>{'<br/>'.join(texts)}</p>")
    return ('\n'.join(slides_html)).encode('utf-8')


def render_preview(files_row) -> tuple:
    """Return (mimetype, bytes, filename) for inline preview if possible.
    If no preview available, return None.
    """
    # determine filename and extension
    meta = files_row.file_metadata or {}
    filename = meta.get('filename') or ''
    if not filename and files_row.storage_url:
        filename = files_row.storage_url.split('/')[-1]

    ext = (filename.split('.')[-1].lower() if '.' in filename else '')

    # download raw bytes
    raw = download_bytes(files_row.storage_url)

    try:
        if ext == 'docx':
            html = _extract_docx_to_html(raw)
            return ('text/html', html, filename)
        if ext in ('xls','xlsx'):
            html = _extract_xlsx_to_html(raw)
            return ('text/html', html, filename)
        if ext in ('ppt','pptx'):
            html = _extract_pptx_to_html(raw)
            return ('text/html', html, filename)
    except Exception:
        # fall through to None -> no preview
        return None

    return None
